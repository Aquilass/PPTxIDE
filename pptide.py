from pptx import Presentation

def execute_python_code(text):
    try:
        exec(text)
    except Exception as e:
        print("An error occurred:", e)


def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    extracted_text = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                shape.text = shape.text.replace("‘", "'").replace("’", "'")
                extracted_text.append(shape.text)

    return extracted_text

if __name__ == "__main__":
    pptx_file = "main.pptx"  # 將路徑替換為你的 PowerPoint 文件的實際路徑
    extracted_text = extract_text_from_pptx(pptx_file)
    exec_str = ""
    for text in extracted_text:
        if text != "":
            exec_str += text
        # text = '\n'.join(extracted_text)
    
    # print(exec_str)
    execute_python_code(exec_str)